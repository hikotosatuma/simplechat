#lambda/index.py
import json
import boto3
import logging
import os
import base64
import io
import re

# ログ設定
logger = logging.getLogger()
logger.setLevel(logging.INFO)

# Bedrockクライアントの初期化
bedrock = boto3.client('bedrock-runtime', region_name='us-east-1')

def remove_markdown_formatting(text):
    """マークダウンフォーマッティングを削除する関数"""
    import re
    
    # **強調** を削除
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    
    # *斜体* を削除
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    
    # ***太字斜体*** を削除
    text = re.sub(r'\*\*\*(.*?)\*\*\*', r'\1', text)
    
    # `コード` を削除
    text = re.sub(r'`(.*?)`', r'\1', text)
    
    # [リンク](url) を削除
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    
    return text

def extract_text_from_pdf(file_content):
    """PDFからテキストを抽出（PyPDF2使用）"""
    try:
        import PyPDF2
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"PDF processing error: {str(e)}")
        return f"PDFの処理中にエラーが発生しました: {str(e)}"

def extract_text_from_ppt_legacy(file_content):
    """古い形式のPPTファイルからテキストを抽出"""
    try:
        import re
        import struct
        
        logger.info("Processing legacy PPT file")
        
        # PPTファイルの基本的なテキスト抽出
        # 複数のエンコーディングで試行し、最も読みやすい結果を返す
        
        extracted_texts = []
        
        # 方法1: UTF-8で試行
        try:
            text_utf8 = file_content.decode('utf-8', errors='ignore')
            # 英数字と日本語文字を抽出
            readable_chars = re.findall(r'[a-zA-Z0-9\s\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FAF.,!?;:\-()\'\"]{4,}', text_utf8)
            utf8_texts = [s.strip() for s in readable_chars if len(s.strip()) >= 4]
            extracted_texts.extend(utf8_texts[:20])  # 最大20個
        except:
            pass
        
        # 方法2: Shift-JIS（日本語）で試行
        try:
            text_sjis = file_content.decode('shift-jis', errors='ignore')
            readable_chars = re.findall(r'[a-zA-Z0-9\s\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FAF.,!?;:\-()\'\"]{4,}', text_sjis)
            sjis_texts = [s.strip() for s in readable_chars if len(s.strip()) >= 4]
            extracted_texts.extend(sjis_texts[:20])
        except:
            pass
        
        # 方法3: バイト列から直接英数字を検索
        try:
            ascii_strings = re.findall(rb'[a-zA-Z0-9\s.,!?;:\-()\'\"]{4,}', file_content)
            ascii_texts = [s.decode('ascii', errors='ignore').strip() for s in ascii_strings]
            ascii_texts = [s for s in ascii_texts if len(s) >= 4 and s.isascii()]
            extracted_texts.extend(ascii_texts[:15])
        except:
            pass
        
        # 重複除去と品質フィルタリング
        unique_texts = []
        seen = set()
        
        for text in extracted_texts:
            clean_text = re.sub(r'\s+', ' ', text.strip())  # 連続する空白を一つに
            
            if (clean_text and 
                len(clean_text) >= 4 and 
                len(clean_text) <= 100 and 
                clean_text not in seen and
                not re.match(r'^[^\w\s]*$', clean_text) and  # 記号のみでない
                not re.match(r'^[\x00-\x1F\x7F-\x9F]+$', clean_text)):  # 制御文字のみでない
                
                unique_texts.append(clean_text)
                seen.add(clean_text)
        
        # 結果を生成
        if unique_texts:
            # 文字数で昇順ソート（短いものから）
            unique_texts.sort(key=len)
            
            result = "--- PPTファイル（レガシー形式）から抽出されたテキスト ---\n\n"
            
            # 最大15個のテキストを表示
            displayed_texts = unique_texts[:15]
            for i, text in enumerate(displayed_texts, 1):
                result += f"{i}. {text}\n"
            
            if len(unique_texts) > 15:
                result += f"\n... 他に{len(unique_texts) - 15}個のテキスト断片があります\n"
            
            result += "\n注意: 古いPPT形式のため、文字化けや不完全な抽出が発生する場合があります。"
            result += "\nより正確な抽出のため、PowerPointで.pptx形式に変換してから再アップロードすることをお勧めします。"
            
            return result
        else:
            return ("PPTファイルからテキストを抽出できませんでした。\n\n"
                   "このファイルは古いPowerPoint形式のため、テキスト抽出が困難です。\n"
                   "解決策: PowerPointでファイルを開き、「名前を付けて保存」で .pptx 形式に変換してから再度アップロードしてください。")
            
    except Exception as e:
        logger.error(f"PPT legacy processing error: {str(e)}")
        return f"古いPPTファイルの処理中にエラーが発生しました: {str(e)}\n\n解決策: PowerPointでファイルを開き、.pptx形式で保存し直してください。"

def extract_text_from_pptx(file_content):
    """PPTXからテキストを抽出（python-pptx使用）"""
    try:
        from pptx import Presentation
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"--- スライド {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
            text += "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"PPTX processing error: {str(e)}")
        return f"PPTXの処理中にエラーが発生しました: {str(e)}"

def extract_text_from_docx(file_content):
    """DOCXからテキストを抽出（python-docx使用）"""
    try:
        from docx import Document
        docx_file = io.BytesIO(file_content)
        doc = Document(docx_file)
        text = ""
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"DOCX processing error: {str(e)}")
        return f"DOCXの処理中にエラーが発生しました: {str(e)}"

def process_file_content(file_content, file_type, file_name):
    """ファイル内容を処理してテキストを抽出"""
    logger.info(f"Processing file: {file_name}, type: {file_type}")
    
    try:
        if file_type == 'application/pdf' or file_name.lower().endswith('.pdf'):
            return extract_text_from_pdf(file_content)
        
        elif (file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' 
              or file_name.lower().endswith('.pptx')):
            return extract_text_from_pptx(file_content)
        
        elif (file_type == 'application/vnd.ms-powerpoint' 
              or file_name.lower().endswith('.ppt')):
            return extract_text_from_ppt_legacy(file_content)
        
        elif (file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
              or file_name.lower().endswith('.docx')):
            return extract_text_from_docx(file_content)
        
        elif (file_type == 'application/msword'
              or file_name.lower().endswith('.doc')):
            return "古いDOCファイル形式です。.docx形式で保存し直してアップロードしてください。"
        
        elif file_type.startswith('text/') or file_name.lower().endswith('.txt'):
            return file_content.decode('utf-8')
        
        else:
            return f"サポートされていないファイル形式です: {file_type}。対応形式: PDF, PPTX, PPT, DOCX, TXT"
    
    except Exception as e:
        logger.error(f"File processing error: {str(e)}")
        return f"ファイル処理エラー: {str(e)}"

def handle_file_upload(event):
    """ファイルアップロードを処理（Lambda内で完結）"""
    try:
        body = json.loads(event.get("body", "{}"))
        file_data = body.get("file")
        file_name = body.get("fileName")
        file_type = body.get("fileType", "")
        
        if not file_data or not file_name:
            return {
                "statusCode": 400,
                "headers": {
                    "Content-Type": "application/json",
                    "Access-Control-Allow-Origin": "*"
                },
                "body": json.dumps({
                    "success": False,
                    "error": "ファイルデータまたはファイル名が見つかりません"
                })
            }
        
        # Base64デコード
        try:
            file_content = base64.b64decode(file_data)
        except Exception as e:
            return {
                "statusCode": 400,
                "headers": {
                    "Content-Type": "application/json",
                    "Access-Control-Allow-Origin": "*"
                },
                "body": json.dumps({
                    "success": False,
                    "error": f"ファイルデータのデコードに失敗しました: {str(e)}"
                })
            }
        
        # ファイルサイズチェック（10MB制限）
        if len(file_content) > 10 * 1024 * 1024:
            return {
                "statusCode": 400,
                "headers": {
                    "Content-Type": "application/json",
                    "Access-Control-Allow-Origin": "*"
                },
                "body": json.dumps({
                    "success": False,
                    "error": "ファイルサイズが大きすぎます（10MB以下にしてください）"
                })
            }
        
        # Lambda内でファイルからテキストを抽出
        extracted_text = process_file_content(file_content, file_type, file_name)
        
        # 抽出されたテキストの長さをチェック
        if len(extracted_text) > 50000:  # 50KB制限
            extracted_text = extracted_text[:50000] + "\n...(テキストが長すぎるため切り詰められました)"
        
        return {
            "statusCode": 200,
            "headers": {
                "Content-Type": "application/json",
                "Access-Control-Allow-Origin": "*"
            },
            "body": json.dumps({
                "success": True,
                "message": "ファイルが正常に処理されました",
                "extracted_text": extracted_text,
                "file_name": file_name,
                "file_size": len(file_content),
                "text_length": len(extracted_text)
            })
        }
        
    except Exception as e:
        logger.error(f"ファイルアップロードエラー: {str(e)}", exc_info=True)
        return {
            "statusCode": 500,
            "headers": {
                "Content-Type": "application/json",
                "Access-Control-Allow-Origin": "*"
            },
            "body": json.dumps({
                "success": False,
                "error": str(e)
            })
        }

def lambda_handler(event, context):
    try:
        logger.info(f"Received event: {json.dumps(event)}")
        
        # HTTPメソッドとパスを取得
        http_method = event.get('httpMethod', '')
        path = event.get('path', '').rstrip('/')
        resource = event.get('resource', '')
        
        logger.info(f"HTTP Method: {http_method}, Path: {path}, Resource: {resource}")
        
        # OPTIONSリクエスト（CORS preflight）への対応
        if http_method == 'OPTIONS':
            return {
                "statusCode": 200,
                "headers": {
                    "Access-Control-Allow-Origin": "*",
                    "Access-Control-Allow-Headers": "Content-Type,X-Amz-Date,Authorization,X-Api-Key,X-Amz-Security-Token",
                    "Access-Control-Allow-Methods": "OPTIONS,POST"
                },
                "body": ""
            }
        
        # パスに基づいてルーティング
        if '/upload' in path or '/upload' in resource:
            return handle_file_upload(event)
        elif '/chat' in path or '/chat' in resource:
            return handle_chat(event)
        else:
            # デフォルトはチャット機能
            return handle_chat(event)
    
    except Exception as e:
        logger.error(f"Error in lambda_handler: {str(e)}", exc_info=True)
        return {
            "statusCode": 500,
            "headers": {
                "Content-Type": "application/json",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Headers": "Content-Type,X-Amz-Date,Authorization,X-Api-Key,X-Amz-Security-Token",
                "Access-Control-Allow-Methods": "OPTIONS,POST"
            },
            "body": json.dumps({
                "success": False,
                "error": str(e)
            })
        }

def handle_chat(event):
    try:
        logger.info("Handling chat request")
        
        # 1) リクエストボディを取得・解析
        body = json.loads(event.get("body", "{}"))
        message = body.get("message", "")
        conversation_history = body.get("conversationHistory", [])
        uploaded_files = body.get("uploadedFiles", [])
        
        logger.info(f"Processing message: {message}")
        
        # 2) アップロードされたファイル情報を含めてコンテキストを構築
        context_message = message
        if uploaded_files:
            context_message = f"アップロードされたファイル情報を参考に回答してください。\n\n"
            for file_info in uploaded_files:
                context_message += f"ファイル名: {file_info.get('name', 'Unknown')}\n"
                context_message += f"抽出テキスト: {file_info.get('extractedText', 'No text')}\n\n"
            context_message += f"ユーザーの質問: {message}"

        # 3) Nova Liteモデル用のリクエストペイロードを作成
        messages = []
        
        # 会話履歴を追加（システムメッセージを除く）
        for msg in conversation_history:
            if msg.get("role") != "system":
                messages.append({
                    "role": msg["role"],
                    "content": [{"text": msg["content"]}]
                })
        
        # 現在のユーザーメッセージを追加
        messages.append({
            "role": "user",
            "content": [{"text": context_message}]
        })

        payload = {
            "messages": messages,
            "inferenceConfig": {
                "temperature": 0.7,
                "topP": 0.9,
                "maxTokens": 1024
            }
        }
        
        logger.info(f"Bedrock payload: {json.dumps(payload, ensure_ascii=False)}")

        # 4) Bedrockでテキスト生成を実行
        response = bedrock.invoke_model(
            modelId='us.amazon.nova-lite-v1:0',
            body=json.dumps(payload),
            contentType='application/json'
        )

        # 5) レスポンスを解析
        response_body = json.loads(response['body'].read())
        logger.info(f"Bedrock response: {json.dumps(response_body, ensure_ascii=False)}")
        
        # Nova Liteのレスポンス形式から生成テキストを取得
        generated_text = response_body['output']['message']['content'][0]['text']
        
        # マークダウン形式を削除（**強調**など）
        generated_text = remove_markdown_formatting(generated_text)

        # 記述問題の採点が要求されているかチェック
        if is_essay_grading_request(message, context_message, uploaded_files):
            essay_score = grade_essay_answer(message, context_message, uploaded_files)
            if essay_score:
                generated_text = essay_score

        # 6) 会話履歴に新しいメッセージを追加
        updated_history = conversation_history + [
            {"role": "user", "content": message},
            {"role": "assistant", "content": generated_text}
        ]

        # 7) 正常レスポンスを返す
        return {
            "statusCode": 200,
            "headers": {
                "Content-Type": "application/json",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Headers": "Content-Type,X-Amz-Date,Authorization,X-Api-Key,X-Amz-Security-Token",
                "Access-Control-Allow-Methods": "OPTIONS,POST"
            },
            "body": json.dumps({
                "success": True,
                "response": generated_text,
                "conversationHistory": updated_history
            }, ensure_ascii=False)
        }

    except Exception as e:
        logger.error(f"Error occurred in chat: {str(e)}", exc_info=True)
        # エラー時のレスポンス
        return {
            "statusCode": 500,
            "headers": {
                "Content-Type": "application/json",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Headers": "Content-Type,X-Amz-Date,Authorization,X-Api-Key,X-Amz-Security-Token",
                "Access-Control-Allow-Methods": "OPTIONS,POST"
            },
            "body": json.dumps({
                "success": False,
                "error": str(e)
            })
        }

def is_essay_grading_request(message, context_message, uploaded_files):
    """記述問題の採点要求かどうかを判定する関数"""
    # 特定のキーワードが含まれているかチェック
    grading_keywords = [
        'GRADE_ESSAY_ANSWER',  # フロントエンドから送信される特別なマーカー
        '記述問題の回答を採点',
        '記述問題採点'
    ]
    
    for keyword in grading_keywords:
        if keyword in message or keyword in context_message:
            return True
    
    return False

def grade_essay_answer(message, context_message, uploaded_files):
    """記述問題の回答を採点する関数"""
    try:
        logger.info("記述問題の採点を開始")
        
        # メッセージから採点に必要な情報を抽出
        import json
        
        # GRADE_ESSAY_ANSWERマーカーの後にJSON形式で情報が含まれていることを想定
        if 'GRADE_ESSAY_ANSWER:' in message:
            json_part = message.split('GRADE_ESSAY_ANSWER:')[1].strip()
            try:
                grading_info = json.loads(json_part)
                user_answer = grading_info.get('userAnswer', '')
                question = grading_info.get('question', {})
                points = grading_info.get('points', [])
                explanation = grading_info.get('explanation', '')
                
                logger.info(f"採点情報: 問題={question.get('question', '')[:50]}, 回答長={len(user_answer)}, ポイント数={len(points)}")
                
                # Bedrockを使用して採点
                grading_prompt = f"""以下の記述問題の回答を採点してください。

【問題】
{question.get('question', '')}

【学生の回答】
{user_answer}

【解答例】
{explanation}

【採点ポイント】
{chr(10).join([f"- {point}" for point in points])}

【採点指示】
1. 各採点ポイントについて、学生の回答が該当するかを評価してください
2. 100点満点で点数をつけてください
3. 具体的なフィードバックを提供してください
4. 改善点があれば指摘してください

以下の形式で回答してください：

得点: XX/100点

各ポイントの評価:
- ポイント1: ○/△/×（理由）
- ポイント2: ○/△/×（理由）
...

総合評価:
（全体的な評価コメント）

改善点:
（具体的な改善提案）"""

                # Bedrockで採点を実行
                messages = [{
                    "role": "user",
                    "content": [{"text": grading_prompt}]
                }]

                payload = {
                    "messages": messages,
                    "inferenceConfig": {
                        "temperature": 0.3,  # 採点では一貫性を重視するため低めに設定
                        "topP": 0.8,
                        "maxTokens": 1024
                    }
                }

                response = bedrock.invoke_model(
                    modelId='us.amazon.nova-lite-v1:0',
                    body=json.dumps(payload),
                    contentType='application/json'
                )

                response_body = json.loads(response['body'].read())
                grading_result = response_body['output']['message']['content'][0]['text']
                
                # マークダウンフォーマッティングを削除
                grading_result = remove_markdown_formatting(grading_result)
                
                logger.info("記述問題の採点が完了")
                return f"📝 記述問題の採点結果\n\n{grading_result}"
                
            except json.JSONDecodeError as e:
                logger.error(f"採点情報のJSON解析エラー: {str(e)}")
                return None
        
        return None
        
    except Exception as e:
        logger.error(f"記述問題採点エラー: {str(e)}", exc_info=True)
        return None
